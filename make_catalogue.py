"""
Generate enterprise catalogue for Tech Strategy Lab (marketing-lab).
Run: python make_catalogue.py
Output: Tech_Strategy_Lab_Catalogue.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os, datetime

# ── Brand colours ──────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x1a, 0x1a, 0x6b)   # #1a1a6b
ORANGE = RGBColor(0xF4, 0x79, 0x20)   # #F47920
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT  = RGBColor(0xF8, 0xF8, 0xF8)
GRAY   = RGBColor(0xAA, 0xAA, 0xAA)
DARK   = RGBColor(0x11, 0x11, 0x11)
GREEN  = RGBColor(0x10, 0xB9, 0x81)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height

BLANK = prs.slide_layouts[6]   # completely blank

# ── Helpers ────────────────────────────────────────────────────────────────────

def _slide():
    return prs.slides.add_slide(BLANK)

def _rect(slide, left, top, width, height, fill, alpha=None):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    return shape

def _txt(slide, text, left, top, width, height,
         size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
         italic=False, wrap=True):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return tb

def _accent_bar(slide, top, color=ORANGE, height=Inches(0.04)):
    _rect(slide, 0, top, W, height, color)

def _screenshot_placeholder(slide, left, top, width, height, label):
    """Draw a labelled placeholder box where a screenshot should be inserted."""
    box = slide.shapes.add_shape(1, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xE8, 0xEC, 0xF5)
    box.line.color.rgb = NAVY
    box.line.width = Pt(1.2)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = f"📸  {label}"
    r.font.size = Pt(10)
    r.font.color.rgb = NAVY
    r.font.italic = True

def _feature_card(slide, left, top, width, height, icon, title, body):
    _rect(slide, left, top, width, height, WHITE)
    border = slide.shapes.add_shape(1, left, top, Inches(0.04), height)
    border.fill.solid(); border.fill.fore_color.rgb = ORANGE; border.line.fill.background()
    _txt(slide, icon, left + Inches(0.12), top + Inches(0.08),
         Inches(0.4), Inches(0.4), size=18, color=DARK)
    _txt(slide, title, left + Inches(0.55), top + Inches(0.10),
         width - Inches(0.65), Inches(0.35), size=11, bold=True, color=DARK)
    _txt(slide, body, left + Inches(0.12), top + Inches(0.48),
         width - Inches(0.22), height - Inches(0.6), size=9, color=RGBColor(0x44,0x44,0x44))

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ══════════════════════════════════════════════════════════════════════════════
sl = _slide()
_rect(sl, 0, 0, W, H, NAVY)
_rect(sl, 0, H - Inches(0.12), W, Inches(0.12), ORANGE)
_rect(sl, W * 0.55, 0, W * 0.45, H, RGBColor(0x14,0x14,0x55))

_txt(sl, "ENTERPRISE CATALOGUE", Inches(0.7), Inches(0.8),
     Inches(5), Inches(0.4), size=9, bold=True, color=ORANGE)
_txt(sl, "Tech Strategy Lab", Inches(0.7), Inches(1.3),
     Inches(6), Inches(1.2), size=40, bold=True, color=WHITE)
_txt(sl, "Marketing Intelligence Platform", Inches(0.7), Inches(2.55),
     Inches(6), Inches(0.5), size=18, color=GRAY)
_txt(sl, (
    "Simulate campaign scenarios, isolate demand shocks,\n"
    "reconstruct brand DNA, and project revenue outcomes\n"
    "with enterprise-grade precision."
), Inches(0.7), Inches(3.3), Inches(5.5), Inches(1.5),
    size=12, color=RGBColor(0xCC,0xCC,0xCC), italic=False)

_txt(sl, "🧬  Powered by Seasonal DNA Engine™",
     Inches(0.7), Inches(5.1), Inches(5), Inches(0.4),
     size=10, bold=True, color=ORANGE)

_txt(sl, "📸  INSERT: App hero screenshot (login or dashboard overview)",
     Inches(7.2), Inches(1.0), Inches(5.8), Inches(5.2),
     size=10, color=RGBColor(0x88,0x88,0xBB), italic=True, align=PP_ALIGN.CENTER)

_txt(sl, f"Confidential  ·  {datetime.date.today().year}",
     Inches(0.7), Inches(6.9), Inches(4), Inches(0.4),
     size=8, color=GRAY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — THE PROBLEM WE SOLVE
# ══════════════════════════════════════════════════════════════════════════════
sl = _slide()
_rect(sl, 0, 0, W, Inches(1.3), NAVY)
_rect(sl, 0, Inches(1.3), W, H - Inches(1.3), LIGHT)
_accent_bar(sl, Inches(1.3))

_txt(sl, "The Problem", Inches(0.7), Inches(0.22),
     Inches(8), Inches(0.6), size=26, bold=True, color=WHITE)
_txt(sl, "Why revenue forecasting fails without a DNA-aware engine",
     Inches(0.7), Inches(0.78), Inches(8), Inches(0.4),
     size=11, color=GRAY)

problems = [
    ("📉", "Seasonal blindness",
     "Standard models treat every week equally. High-traffic months are understated, "
     "slow periods are overestimated — leading to budget misallocation."),
    ("🎯", "Campaign noise contamination",
     "A single flash-sale week corrupts the entire demand baseline. Teams cannot "
     "separate organic demand from artificial spikes."),
    ("🔮", "No 'what-if' scenario engine",
     "Decision-makers lack a safe sandbox to model 'what happens if we replicate "
     "last year's Black Friday push in March?' before committing spend."),
    ("🧩", "Brand comparison impossible",
     "Multi-brand portfolios carry incomparable seasonality profiles. Aggregating "
     "across brands without DNA normalisation produces meaningless KPIs."),
]
left_col = [0, 2]
for i, (icon, title, body) in enumerate(problems):
    col = i % 2
    row = i // 2
    _feature_card(
        sl,
        left=Inches(0.5 + col * 6.3),
        top=Inches(1.7 + row * 2.4),
        width=Inches(6.0),
        height=Inches(2.1),
        icon=icon, title=title, body=body,
    )

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — HOW IT WORKS (DNA ENGINE)
# ══════════════════════════════════════════════════════════════════════════════
sl = _slide()
_rect(sl, 0, 0, W, Inches(1.3), NAVY)
_rect(sl, 0, Inches(1.3), W, H - Inches(1.3), WHITE)
_accent_bar(sl, Inches(1.3))

_txt(sl, "How It Works", Inches(0.7), Inches(0.22),
     Inches(8), Inches(0.6), size=26, bold=True, color=WHITE)
_txt(sl, "The Seasonal DNA Engine™ — a three-layer demand reconstruction model",
     Inches(0.7), Inches(0.78), Inches(9), Inches(0.4), size=11, color=GRAY)

steps = [
    ("1", "Historical DNA Extraction",
     "Upload brand data. The engine computes daily, weekly, and monthly seasonality "
     "indices (clicks, CR, AOV) across all available years — weighted by recency and similarity."),
    ("2", "Trial Calibration",
     "Enter your trial-period actuals (clicks, orders, revenue). The engine strips "
     "adjustments (promotions, suppressions) to derive clean base constants."),
    ("3", "Forward Projection",
     "DNA indices are multiplied by base constants to produce daily Baseline and "
     "Simulation curves with ±15 % confidence bands for the full projection year."),
    ("4", "Event Simulation",
     "Inject campaigns, DNA swaps, de-shocked signatures, or custom drag events. "
     "Each event is layered on top — fully attributed and reversible."),
]
for i, (num, title, body) in enumerate(steps):
    top = Inches(1.55 + i * 1.42)
    _rect(sl, Inches(0.5), top + Inches(0.05),
          Inches(0.55), Inches(0.55), ORANGE)
    _txt(sl, num, Inches(0.5), top + Inches(0.05),
         Inches(0.55), Inches(0.55), size=16, bold=True,
         color=WHITE, align=PP_ALIGN.CENTER)
    _txt(sl, title, Inches(1.25), top,
         Inches(4.5), Inches(0.4), size=12, bold=True, color=DARK)
    _txt(sl, body, Inches(1.25), top + Inches(0.42),
         Inches(4.5), Inches(0.85), size=9.5, color=RGBColor(0x44,0x44,0x44))
    if i < 3:
        _rect(sl, Inches(0.74), top + Inches(0.6),
              Inches(0.06), Inches(0.85), GRAY)

_screenshot_placeholder(
    sl, Inches(6.6), Inches(1.6), Inches(6.5), Inches(5.6),
    "INSERT: DNA Drag / Simulation chart screenshot")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — DASHBOARD
# ══════════════════════════════════════════════════════════════════════════════
sl = _slide()
_rect(sl, 0, 0, W, H, LIGHT)
_rect(sl, 0, 0, W, Inches(0.6), NAVY)
_txt(sl, "📊  Main Dashboard", Inches(0.5), Inches(0.1),
     Inches(8), Inches(0.4), size=14, bold=True, color=WHITE)

_screenshot_placeholder(
    sl, Inches(0.3), Inches(0.75), Inches(8.5), Inches(5.5),
    "INSERT: Full Dashboard screenshot (Goal Tracker + DNA chart + Projection chart)")

bullets = [
    "Baseline vs Simulation projection (daily/weekly/monthly)",
    "Goal Tracker — set revenue / orders / clicks targets",
    "Historical DNA comparison across brands",
    "±15 % confidence band visualisation",
]
for i, b in enumerate(bullets):
    _txt(sl, f"✓  {b}", Inches(9.1), Inches(1.0 + i * 0.8),
         Inches(3.9), Inches(0.6), size=10, color=DARK)

_rect(sl, Inches(9.0), Inches(4.5), Inches(4.1), Inches(1.6), NAVY)
_txt(sl, "Used by", Inches(9.1), Inches(4.6), Inches(2), Inches(0.4),
     size=9, color=GRAY)
_txt(sl, "Strategy & Growth teams\nDigital Marketing Directors\nBrand Portfolio Managers",
     Inches(9.1), Inches(4.95), Inches(3.9), Inches(1.0), size=10, color=WHITE)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — SIMULATION LAB
# ══════════════════════════════════════════════════════════════════════════════
sl = _slide()
_rect(sl, 0, 0, W, H, LIGHT)
_rect(sl, 0, 0, W, Inches(0.6), NAVY)
_txt(sl, "⚡  Simulation Lab", Inches(0.5), Inches(0.1),
     Inches(8), Inches(0.4), size=14, bold=True, color=WHITE)

_screenshot_placeholder(
    sl, Inches(0.3), Inches(0.75), Inches(8.5), Inches(5.5),
    "INSERT: Simulation Lab screenshot (Events tab or De-Shock tab)")

features = [
    ("📣", "Campaign Injection",  "Front-Loaded, Linear Fade, Step, Delayed Peak shapes"),
    ("🖱️", "DNA Sculpting",       "Drag-to-reshape seasonal indices month-by-month"),
    ("🧹", "De-Shock Tool",       "Isolate artificial spikes (historical or forecast)"),
    ("🎛️", "Spike Compressor",    "Re-scale any spike — 'what if 30 % less lift?'"),
    ("💉", "Signature Library",   "Save, re-inject, and compare event signatures"),
    ("📋", "Attribution Engine",  "Gap attribution — % of target each event fills"),
]
for i, (ic, tt, dd) in enumerate(features):
    col = i % 2; row = i // 2
    top = Inches(0.9 + row * 1.5)
    lft = Inches(9.15 + col * 1.9)
    _txt(sl, ic, lft, top, Inches(0.4), Inches(0.4), size=14, color=DARK)
    _txt(sl, tt, lft + Inches(0.42), top,
         Inches(1.4), Inches(0.35), size=9, bold=True, color=DARK)
    _txt(sl, dd, lft + Inches(0.42), top + Inches(0.38),
         Inches(1.4), Inches(0.9), size=7.5, color=RGBColor(0x55,0x55,0x55))

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — BRAND FORGE & USER LOG
# ══════════════════════════════════════════════════════════════════════════════
sl = _slide()
_rect(sl, 0, 0, W, H, LIGHT)
_rect(sl, 0, 0, W, Inches(0.6), NAVY)
_txt(sl, "🔬  Brand Forge  ·  📋  User Log", Inches(0.5), Inches(0.1),
     Inches(8), Inches(0.4), size=14, bold=True, color=WHITE)

_screenshot_placeholder(
    sl, Inches(0.3), Inches(0.75), Inches(6.2), Inches(3.1),
    "INSERT: Brand Forge screenshot (DNA preview chart)")
_screenshot_placeholder(
    sl, Inches(0.3), Inches(4.0), Inches(6.2), Inches(3.1),
    "INSERT: User Log screenshot (activity log view)")

forge_points = [
    "Blend two existing brand DNA profiles",
    "Adjust monthly click, CR, and AOV indices",
    "Generate synthetic daily data with log-normal noise",
    "Instantly appears in DNA Brands selector",
]
log_points = [
    "Every click and modification is timestamped",
    "Filter by user, action type, date",
    "Download full audit CSV",
    "Granular login snapshots (data state at login)",
]

_rect(sl, Inches(6.8), Inches(0.75), Inches(6.2), Inches(2.8), NAVY)
_txt(sl, "🔬  Brand Forge", Inches(7.0), Inches(0.9),
     Inches(5.8), Inches(0.4), size=13, bold=True, color=WHITE)
for i, p in enumerate(forge_points):
    _txt(sl, f"→  {p}", Inches(7.0), Inches(1.35 + i * 0.52),
         Inches(5.8), Inches(0.45), size=10, color=GRAY)

_rect(sl, Inches(6.8), Inches(3.75), Inches(6.2), Inches(3.4), RGBColor(0x0D,0x0D,0x44))
_txt(sl, "📋  User Log & Audit Trail", Inches(7.0), Inches(3.9),
     Inches(5.8), Inches(0.4), size=13, bold=True, color=WHITE)
for i, p in enumerate(log_points):
    _txt(sl, f"→  {p}", Inches(7.0), Inches(4.38 + i * 0.55),
         Inches(5.8), Inches(0.45), size=10, color=GRAY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — ENTERPRISE FEATURES
# ══════════════════════════════════════════════════════════════════════════════
sl = _slide()
_rect(sl, 0, 0, W, Inches(1.3), NAVY)
_rect(sl, 0, Inches(1.3), W, H - Inches(1.3), WHITE)
_accent_bar(sl, Inches(1.3))

_txt(sl, "Enterprise-Grade Capabilities", Inches(0.7), Inches(0.22),
     Inches(10), Inches(0.6), size=26, bold=True, color=WHITE)
_txt(sl, "Built for multi-brand portfolio management at scale",
     Inches(0.7), Inches(0.78), Inches(9), Inches(0.4), size=11, color=GRAY)

cards = [
    ("🔐", "Role-Based Access",
     "Single-tenant deployment with username/password auth. Full audit trail of every action by every user."),
    ("📊", "Multi-Brand DNA",
     "Support for unlimited brands. Similarity-weighted DNA blending across years and entities."),
    ("📁", "Excel Strategy Reports",
     "One-click download of full projection data, event log, and DNA weights as formatted Excel."),
    ("🌍", "Multilingual UI",
     "English and Italian interface. Extend to any language via the i18n configuration layer."),
    ("🔄", "Live Data Ingestion",
     "Upload CSV to add or replace brand data at any time. Brand Forge creates synthetic brands from DNA."),
    ("⚙️", "Configurable Defaults",
     "Per-brand campaign impact coefficients. Global defaults with brand-level overrides. Settings saved persistently."),
]
for i, (ic, tt, dd) in enumerate(cards):
    col = i % 3; row = i // 3
    _feature_card(
        sl,
        left=Inches(0.45 + col * 4.28),
        top=Inches(1.55 + row * 2.55),
        width=Inches(3.95),
        height=Inches(2.25),
        icon=ic, title=tt, body=dd,
    )

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — TECHNICAL STACK
# ══════════════════════════════════════════════════════════════════════════════
sl = _slide()
_rect(sl, 0, 0, W, H, LIGHT)
_rect(sl, 0, 0, Inches(4.5), H, NAVY)

_txt(sl, "Technical\nStack", Inches(0.5), Inches(0.8),
     Inches(3.5), Inches(1.2), size=28, bold=True, color=WHITE)
_txt(sl, (
    "Python-native, fully open-source stack.\n"
    "Runs on-premise or any cloud provider.\n"
    "No external data dependencies."
), Inches(0.5), Inches(2.2), Inches(3.5), Inches(2.0),
    size=11, color=GRAY, italic=False)

_txt(sl, "Deploy on:", Inches(0.5), Inches(4.5),
     Inches(3.5), Inches(0.35), size=10, bold=True, color=ORANGE)
_txt(sl, "Streamlit Community Cloud  ·  AWS  ·  GCP\nAzure  ·  On-premise Docker",
     Inches(0.5), Inches(4.9), Inches(3.5), Inches(0.8), size=10, color=GRAY)

stack = [
    ("Frontend", "Streamlit 1.32+", "Interactive web app, no JavaScript required"),
    ("Data Layer", "Pandas 2.0 + NumPy 1.26", "Fast in-memory computation"),
    ("Visualisation", "Plotly 5.20", "Interactive charts, hover, zoom"),
    ("DNA Engine", "Custom Python modules", "engine/dna.py, calibration.py, simulation.py"),
    ("Persistence", "CSV flat files", "profiles, dataset, settings, activity log"),
    ("Export", "OpenPyXL 3.1", "Formatted Excel reports with styling"),
    ("Auth", "Session-state auth gate", "Username + password, extensible to SSO"),
]
for i, (layer, tech, note) in enumerate(stack):
    top = Inches(0.75 + i * 0.85)
    _rect(sl, Inches(4.8), top, Inches(8.2), Inches(0.7), WHITE)
    _txt(sl, layer, Inches(5.0), top + Inches(0.05),
         Inches(1.8), Inches(0.3), size=8, bold=True, color=GRAY)
    _txt(sl, tech, Inches(5.0), top + Inches(0.32),
         Inches(3.5), Inches(0.3), size=11, bold=True, color=DARK)
    _txt(sl, note, Inches(8.7), top + Inches(0.22),
         Inches(4.0), Inches(0.4), size=9, color=RGBColor(0x66,0x66,0x66))

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — CLOSING / CONTACT
# ══════════════════════════════════════════════════════════════════════════════
sl = _slide()
_rect(sl, 0, 0, W, H, NAVY)
_rect(sl, 0, H - Inches(0.12), W, Inches(0.12), ORANGE)

_txt(sl, "Ready to deploy?", Inches(1.5), Inches(1.2),
     Inches(10), Inches(1.1), size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
_txt(sl, (
    "Tech Strategy Lab is available as a private cloud deployment\n"
    "or on-premise installation with full source-code access."
), Inches(1.5), Inches(2.5), Inches(10), Inches(0.9),
    size=14, color=GRAY, align=PP_ALIGN.CENTER)

_rect(sl, Inches(4.0), Inches(3.7), Inches(5.33), Inches(1.8), RGBColor(0x14,0x14,0x55))
_txt(sl, "📧  parsa.hajiannejad@universitybox.com",
     Inches(4.2), Inches(3.9), Inches(5.0), Inches(0.5),
     size=12, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
_txt(sl, "🔗  github.com/Parsa-Hajian/marketing-lab",
     Inches(4.2), Inches(4.45), Inches(5.0), Inches(0.5),
     size=11, color=WHITE, align=PP_ALIGN.CENTER)

_txt(sl, "🧬  Tech Strategy Lab  ·  Marketing Intelligence Platform",
     Inches(1.5), Inches(6.0), Inches(10.3), Inches(0.5),
     size=10, color=GRAY, align=PP_ALIGN.CENTER)
_txt(sl, f"Confidential  ·  {datetime.date.today().year}",
     Inches(1.5), Inches(6.9), Inches(10.3), Inches(0.35),
     size=8, color=RGBColor(0x55,0x55,0x77), align=PP_ALIGN.CENTER)

# ── Save ───────────────────────────────────────────────────────────────────────
out = "Tech_Strategy_Lab_Catalogue.pptx"
prs.save(out)
print(f"✅  Saved: {out}")
